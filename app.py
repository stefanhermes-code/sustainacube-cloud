import streamlit as st
import openai
import os
from pathlib import Path
import PyPDF2
from docx import Document
from dotenv import load_dotenv
import re
from collections import Counter
import json
import time
from datetime import datetime

# Load environment variables
load_dotenv()

class SustainaCubeMinimal:
    def __init__(self):
        # Try Streamlit secrets first, fallback to environment variables
        try:
            api_key = st.secrets["OPENAI_API_KEY"]
            assistant_id = st.secrets.get("OPENAI_ASSISTANT_ID", "")
        except:
            # Fallback to environment variables for local development
            api_key = os.getenv("OPENAI_API_KEY")
            assistant_id = os.getenv("OPENAI_ASSISTANT_ID", "")
        
        self.openai_client = openai.OpenAI(api_key=api_key)
        self.documents = []
        self.processed = False
        # Assistant (OpenAI Assistants API)
        self.assistant_id = assistant_id
        # Resume/cache/tracking
        self.cache_path = Path(__file__).parent / "processed_cache.json"
        self.chunks_store_path = Path(__file__).parent / "chunks_store.jsonl"
        self.progress_path = Path(__file__).parent / ".progress.json"
        self.log_file = Path(__file__).parent / "processing_log.txt"
        self.log_entries = []
        self.lock_file = Path(__file__).parent / "app_lock.json"
        self._processed_index = {}  # file_path -> {mtime, size}
        self._load_cache_and_chunks()
        
    def extract_text_from_file(self, file_path):
        """Extract text from various file formats"""
        try:
            file_path = Path(file_path)
            if file_path.suffix.lower() == '.pdf':
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                    
                    # If no text extracted, try OCR
                    if not text.strip():
                        try:
                            import pytesseract
                            from PIL import Image
                            import fitz  # PyMuPDF for better PDF handling
                            
                            doc = fitz.open(file_path)
                            ocr_text = ""
                            for page_num in range(len(doc)):
                                page = doc.load_page(page_num)
                                pix = page.get_pixmap()
                                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                                ocr_text += pytesseract.image_to_string(img) + "\n"
                            doc.close()
                            return ocr_text if ocr_text.strip() else ""
                        except ImportError:
                            return f"[PDF file: {file_path.name} - OCR libraries not installed]"
                        except Exception as e:
                            return f"[PDF file: {file_path.name} - OCR failed: {str(e)}]"
                    
                    return text
            elif file_path.suffix.lower() == '.docx':
                doc = Document(file_path)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                return text
            elif file_path.suffix.lower() == '.txt':
                with open(file_path, 'r', encoding='utf-8') as file:
                    return file.read()
            elif file_path.suffix.lower() == '.csv':
                # Read CSV as plain text (simple, dependency-free)
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            elif file_path.suffix.lower() in ['.md']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            elif file_path.suffix.lower() in ['.html', '.htm']:
                # Naive HTML to text: strip tags
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    raw = f.read()
                text = re.sub(r'<script[\s\S]*?</script>', ' ', raw, flags=re.IGNORECASE)
                text = re.sub(r'<style[\s\S]*?</style>', ' ', text, flags=re.IGNORECASE)
                text = re.sub(r'<[^>]+>', ' ', text)
                text = re.sub(r'\s+', ' ', text)
                return text
            elif file_path.suffix.lower() in ['.pptx', '.ppt']:
                # Extract text from PowerPoint files
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    return text
                except ImportError:
                    return f"[PowerPoint file: {file_path.name} - python-pptx not installed]"
            elif file_path.suffix.lower() in ['.xlsx', '.xls']:
                # Extract text from Excel files
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(file_path)
                    text = ""
                    for sheet_name in wb.sheetnames:
                        sheet = wb[sheet_name]
                        for row in sheet.iter_rows(values_only=True):
                            row_text = " ".join(str(cell) for cell in row if cell is not None)
                            if row_text.strip():
                                text += row_text + "\n"
                    return text
                except ImportError:
                    return f"[Excel file: {file_path.name} - openpyxl not installed]"
            elif file_path.suffix.lower() == '.doc':
                # Extract text from legacy Word files
                try:
                    import docx2txt
                    return docx2txt.process(file_path)
                except ImportError:
                    return f"[Legacy Word file: {file_path.name} - python-docx2txt not installed]"
            else:
                return ""
        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
            return ""
    
    def chunk_text(self, text, chunk_size=1000, overlap=200):
        """Split text into overlapping chunks for better retrieval"""
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            chunks.append(chunk)
            start = end - overlap
        return chunks
    
    def process_documents(self, documents_folder):
        """Process all documents in the folder with resume, ETA, and progress tracking"""
        st.write("🔄 Processing documents...")
        start_time = time.time()
        
        # Initialize logging
        self._log("info", f"Starting document processing for: {documents_folder}")
        failed_files = {}
        skipped_files = {}
        successful_files = []
        processed_count = 0
        
        # Discover candidate files recursively (exclude program dir)
        base_path = Path(documents_folder)
        program_dir_name = 'SustainaCube_RAG'
        patterns = ['**/*.pdf', '**/*.docx', '**/*.txt', '**/*.csv', '**/*.md', '**/*.html', '**/*.htm', '**/*.pptx', '**/*.xlsx', '**/*.xls', '**/*.ppt', '**/*.doc']
        all_files = []
        for pattern in patterns:
            all_files.extend(base_path.rglob(pattern))
        # Skip temp/lock files like ~$*.docx and anything under program dir
        doc_files = []
        for p in all_files:
            if p.name.startswith('~$'):
                continue
            parts = set(p.parts)
            if program_dir_name in parts:
                continue
            doc_files.append(p)
        total_files = len(doc_files)
        processed = 0
        
        # Load prior documents so resume works
        existing_docs = len(self.documents)
        progress_bar = st.progress(0.0)
        current_file_placeholder = st.empty()
        eta_placeholder = st.empty()
        count_placeholder = st.empty()
        
        for file_path in doc_files:
            try:
                stat = file_path.stat()
                key = str(file_path.resolve())
                mtime = int(stat.st_mtime)
                size = stat.st_size
                already = self._processed_index.get(key)
                # If unchanged and already processed, skip
                if already and already.get('mtime') == mtime and already.get('size') == size:
                    processed += 1
                    self._write_progress(processed, total_files, file_path.name, start_time)
                    self._update_progress_ui(progress_bar, current_file_placeholder, eta_placeholder, count_placeholder, processed, total_files, file_path.name, start_time)
                    continue
            except Exception:
                pass
            
            current_file_placeholder.write(f"Processing: {file_path.name}")
            
            # Extract text
            try:
                text = self.extract_text_from_file(file_path)
                if not text.strip():
                    self._log("warning", "No text extracted from file", str(file_path))
                    skipped_files[str(file_path)] = "No text content"
                    processed += 1
                    self._write_progress(processed, total_files, file_path.name, start_time)
                    self._update_progress_ui(progress_bar, current_file_placeholder, eta_placeholder, count_placeholder, processed, total_files, file_path.name, start_time)
                    continue
            except Exception as e:
                self._log("error", f"Failed to extract text: {str(e)}", str(file_path))
                failed_files[str(file_path)] = str(e)
                processed += 1
                self._write_progress(processed, total_files, file_path.name, start_time)
                self._update_progress_ui(progress_bar, current_file_placeholder, eta_placeholder, count_placeholder, processed, total_files, file_path.name, start_time)
                continue
            
            # Remove any previously stored segments for this file in memory
            self._remove_existing_file_chunks_in_memory(str(file_path))
            
            # Create document segments
            chunks = self.chunk_text(text)
            
            # Add to documents list and persist to chunks_store
            self._append_chunks(file_path, chunks)
            
            # Log successful processing
            self._log("info", f"Successfully processed {len(chunks)} segments", str(file_path))
            successful_files.append(str(file_path))
            processed_count += 1
            
            # Update cache index for resume
            self._processed_index[str(file_path.resolve())] = {'mtime': int(file_path.stat().st_mtime), 'size': file_path.stat().st_size}
            self._save_cache()
            
            processed += 1
            self._write_progress(processed, total_files, file_path.name, start_time)
            self._update_progress_ui(progress_bar, current_file_placeholder, eta_placeholder, count_placeholder, processed, total_files, file_path.name, start_time)
        
        self.processed = True
        self._finalize_progress()
        
        # Write processing summary to log
        self._write_log_summary(total_files, processed_count, failed_files, skipped_files, successful_files)
        self._log("info", f"Processing completed. Successfully processed: {processed_count}, Failed: {len(failed_files)}, Skipped: {len(skipped_files)}")
        
        st.success(f"✅ Processed {processed} of {total_files} files; total documents loaded: {len(self.documents)} (prev: {existing_docs})")
    
    def search_documents(self, query, n_results=5):
        """Search for relevant document segments using keyword matching and scoring"""
        if not self.processed:
            return []
        
        query_words = set(query.lower().split())
        results = []
        
        for doc in self.documents:
            doc_words = set(doc['words'])
            
            # Calculate relevance score
            common_words = query_words.intersection(doc_words)
            if len(common_words) > 0:
                # Score based on word overlap and frequency
                word_freq = Counter(doc['words'])
                score = sum(word_freq[word] for word in common_words) / len(doc['words'])
                
                results.append({
                    'text': doc['text'],
                    'filename': doc['filename'],
                    'file_path': doc['file_path'],
                    'similarity': score,
                    'matched_words': list(common_words)
                })
        
        # Sort by relevance score and return top results
        results.sort(key=lambda x: x['similarity'], reverse=True)
        return results[:n_results]
    
    def generate_answer(self, query, context_chunks):
        """Generate answer using OpenAI with retrieved context"""
        context = "\n\n".join([chunk['text'] for chunk in context_chunks])
        
        prompt = f"""
        You are a sustainability expert with access to a comprehensive knowledge base of research documents.
        
        Question: {query}
        
        Context from research documents:
        {context}
        
        Please provide a comprehensive, well-structured answer based on the provided context. 
        Include specific details, data points, and cite the source documents when possible.
        If the context doesn't contain enough information to fully answer the question, 
        clearly state what information is missing.
        
        Format your response as a professional sustainability report with:
        1. Executive Summary
        2. Key Findings
        3. Supporting Evidence
        4. Source References
        """
        
        try:
            response = self.openai_client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are a sustainability expert providing detailed, accurate information based on research documents."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=2000,
                temperature=0.3
            )
            return response.choices[0].message.content
        except Exception as e:
            return f"Error generating answer: {e}"
    
    def answer_question(self, question):
        """Main method to answer a question using Assistant if configured, else local retrieval"""
        # Prefer Assistant (Vector Store + WebSearch) when configured
        if self.assistant_id:
            try:
                thread = self.openai_client.beta.threads.create()
                self.openai_client.beta.threads.messages.create(
                    thread_id=thread.id,
                    role="user",
                    content=question
                )
                run = self.openai_client.beta.threads.runs.create(
                    thread_id=thread.id,
                    assistant_id=self.assistant_id
                )
                # Poll until completion
                start = time.time()
                while True:
                    r = self.openai_client.beta.threads.runs.retrieve(thread_id=thread.id, run_id=run.id)
                    if r.status in ["completed", "failed", "cancelled", "expired"]:
                        break
                    time.sleep(0.4)
                    # optional safety timeout
                    if time.time() - start > 120:
                        return "Assistant timeout. Please try again.", []
                if r.status != "completed":
                    return f"Assistant run status: {r.status}", []
                msgs = self.openai_client.beta.threads.messages.list(thread_id=thread.id)
                answer_text = ""
                sources = []
                for m in msgs.data:
                    if m.role == 'assistant':
                        parts = []
                        for c in m.content:
                            if getattr(c, 'type', '') == 'text':
                                parts.append(c.text.value)
                        if parts:
                            answer_text = "\n\n".join(parts)
                        # If the assistant returns file references as attachments
                        if hasattr(m, 'attachments') and m.attachments:
                            for att in m.attachments:
                                fname = getattr(att, 'filename', None)
                                if fname:
                                    sources.append({'filename': fname, 'similarity_score': 1.0})
                        break
                
                # Extract source references from the answer text if no attachments found
                if not sources and answer_text:
                    import re
                    # Look for file references in the text (common patterns)
                    file_refs = re.findall(r'\[([^\]]+\.(?:pdf|docx?|txt|pptx?|xlsx?|html?|md|csv))\]', answer_text, re.IGNORECASE)
                    for ref in file_refs:
                        sources.append({'filename': ref, 'similarity_score': 1.0})
                    
                    # Also look for quoted filenames
                    quoted_files = re.findall(r'"([^"]+\.(?:pdf|docx?|txt|pptx?|xlsx?|html?|md|csv))"', answer_text, re.IGNORECASE)
                    for ref in quoted_files:
                        sources.append({'filename': ref, 'similarity_score': 1.0})
                
                return (answer_text or ""), sources
            except Exception as e:
                return f"Assistant error: {e}", []
        
        # Fallback to local retrieval
        search_results = self.search_documents(question)
        if not search_results:
            return "No relevant information found in the knowledge base.", []
        answer = self.generate_answer(question, search_results)
        sources = []
        for result in search_results:
            sources.append({'filename': result['filename'], 'similarity_score': result['similarity'], 'matched_words': result['matched_words']})
        return answer, sources

    # -------------- Internal helpers for cache/progress --------------
    def _load_cache_and_chunks(self):
        try:
            if self.cache_path.exists():
                with open(self.cache_path, 'r', encoding='utf-8') as f:
                    self._processed_index = json.load(f)
            else:
                self._processed_index = {}
        except Exception:
            self._processed_index = {}
        
        # Load previously stored document segments into memory
        try:
            if self.chunks_store_path.exists():
                with open(self.chunks_store_path, 'r', encoding='utf-8') as f:
                    for line in f:
                        try:
                            rec = json.loads(line)
                            self.documents.append(rec)
                        except Exception:
                            continue
                if self.documents:
                    self.processed = True
        except Exception:
            pass
    
    def _save_cache(self):
        try:
            with open(self.cache_path, 'w', encoding='utf-8') as f:
                json.dump(self._processed_index, f)
        except Exception:
            pass
    
    def _append_chunks(self, file_path: Path, chunks):
        records = []
        for i, chunk in enumerate(chunks):
            if not chunk.strip():
                continue
            rec = {
                'filename': file_path.name,
                'file_path': str(file_path),
                'chunk_id': i,
                'text': chunk,
                'file_type': file_path.suffix,
                'words': chunk.lower().split()
            }
            self.documents.append(rec)
            records.append(rec)
        # Append to persistent store
        try:
            with open(self.chunks_store_path, 'a', encoding='utf-8') as f:
                for rec in records:
                    f.write(json.dumps(rec, ensure_ascii=False) + "\n")
        except Exception:
            pass
    
    def _remove_existing_file_chunks_in_memory(self, file_path_str: str):
        if not self.documents:
            return
        self.documents = [d for d in self.documents if d.get('file_path') != file_path_str]
        # Rebuild chunks_store without this file (simple rewrite)
        try:
            if self.chunks_store_path.exists():
                with open(self.chunks_store_path, 'r', encoding='utf-8') as f:
                    lines = f.readlines()
                with open(self.chunks_store_path, 'w', encoding='utf-8') as f:
                    for line in lines:
                        try:
                            rec = json.loads(line)
                            if rec.get('file_path') != file_path_str:
                                f.write(line)
                        except Exception:
                            continue
        except Exception:
            pass
    
    def _write_progress(self, processed, total, current_file, start_time):
        try:
            elapsed = max(0.001, time.time() - start_time)
            pct = (processed / total) if total else 1.0
            # naive ETA
            eta_sec = (elapsed / processed) * (total - processed) if processed > 0 and total > 0 else 0
            data = {
                'processed': processed,
                'total': total,
                'percent': round(pct * 100, 2),
                'current_file': current_file,
                'elapsed_seconds': int(elapsed),
                'eta_seconds': int(eta_sec)
            }
            with open(self.progress_path, 'w', encoding='utf-8') as f:
                json.dump(data, f)
        except Exception:
            pass
    
    def _update_progress_ui(self, bar, file_ph, eta_ph, count_ph, processed, total, current_file, start_time):
        pct = (processed / total) if total else 1.0
        bar.progress(pct)
        elapsed = max(0.001, time.time() - start_time)
        eta_sec = int((elapsed / processed) * (total - processed)) if processed > 0 and total > 0 else 0
        mins, secs = divmod(max(0, eta_sec), 60)
        eta_ph.write(f"ETA: {mins}m {secs}s remaining")
        count_ph.write(f"Files: {processed}/{total} | Current: {current_file}")
    
    def _finalize_progress(self):
        try:
            if self.progress_path.exists():
                with open(self.progress_path, 'w', encoding='utf-8') as f:
                    json.dump({'status': 'complete'}, f)
        except Exception:
            pass
    
    def _get_document_files(self, doc_dir: str) -> list:
        """Get list of all document files to process"""
        if not os.path.exists(doc_dir):
            return []
        
        doc_files = []
        for ext in ['.pdf', '.docx', '.txt', '.csv', '.md', '.html', '.htm', '.pptx', '.xlsx', '.xls', '.ppt', '.doc']:
            for file_path in Path(doc_dir).rglob(f'*{ext}'):
                # Skip files in SustainaCube_RAG directory and temporary files
                if 'SustainaCube_RAG' not in str(file_path) and not file_path.name.startswith('~$'):
                    doc_files.append(str(file_path))
        return doc_files
    
    def check_for_new_files(self, doc_dir: str) -> int:
        """Check for new or updated files without processing them"""
        if not os.path.exists(doc_dir):
            return 0
            
        # Load existing cache
        cache = self._load_cache_and_chunks()
        
        # If no cache exists, all files are new
        if cache is None:
            return len(self._get_document_files(doc_dir))
        
        new_files = 0
        for file_path in self._get_document_files(doc_dir):
            if file_path not in cache:
                new_files += 1
            else:
                # Check if file was modified
                try:
                    current_mtime = os.path.getmtime(file_path)
                    current_size = os.path.getsize(file_path)
                    if (current_mtime != cache[file_path]['mtime'] or 
                        current_size != cache[file_path]['size']):
                        new_files += 1
                except OSError:
                    new_files += 1
                    
        return new_files
    
    def _log(self, level, message, file_path=None):
        """Add entry to processing log"""
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        entry = {
            "timestamp": timestamp,
            "level": level,
            "message": message,
            "file_path": file_path
        }
        self.log_entries.append(entry)
        
        # Write to log file
        with open(self.log_file, "a", encoding="utf-8") as f:
            f.write(f"[{timestamp}] {level.upper()}: {message}")
            if file_path:
                f.write(f" | File: {file_path}")
            f.write("\n")
    
    def _write_log_summary(self, total_files, processed_files, failed_files, skipped_files, successful_files=None):
        """Write summary to log file"""
        with open(self.log_file, "a", encoding="utf-8") as f:
            f.write(f"\n{'='*60}\n")
            f.write(f"PROCESSING SUMMARY - {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write(f"{'='*60}\n")
            f.write(f"Total files found: {total_files}\n")
            f.write(f"Successfully processed: {processed_files}\n")
            f.write(f"Failed to process: {len(failed_files)}\n")
            f.write(f"Skipped (duplicates/unchanged): {len(skipped_files)}\n")
            
            if successful_files:
                f.write(f"\nSUCCESSFULLY PROCESSED FILES:\n")
                for file_path in successful_files:
                    f.write(f"  - {file_path}\n")
            
            f.write(f"\nFAILED FILES:\n")
            for file_path, reason in failed_files.items():
                f.write(f"  - {file_path}: {reason}\n")
            f.write(f"\nSKIPPED FILES:\n")
            for file_path, reason in skipped_files.items():
                f.write(f"  - {file_path}: {reason}\n")
            f.write(f"{'='*60}\n\n")
    
    def _acquire_lock(self, user_name):
        """Acquire application lock to prevent conflicts"""
        try:
            if self.lock_file.exists():
                with open(self.lock_file, 'r', encoding='utf-8') as f:
                    lock_data = json.load(f)
                # Check if lock is stale (older than 5 minutes)
                if time.time() - lock_data.get('timestamp', 0) > 300:
                    self._release_lock()
                else:
                    return False, lock_data.get('user', 'Unknown')
            
            # Create new lock
            lock_data = {
                'user': user_name,
                'timestamp': time.time(),
                'pid': os.getpid()
            }
            with open(self.lock_file, 'w', encoding='utf-8') as f:
                json.dump(lock_data, f)
            return True, user_name
        except Exception:
            return False, "Error"
    
    def _release_lock(self):
        """Release application lock"""
        try:
            if self.lock_file.exists():
                self.lock_file.unlink()
        except Exception:
            pass
    
    def _check_who_online(self):
        """Check who is currently using the application"""
        try:
            if self.lock_file.exists():
                with open(self.lock_file, 'r', encoding='utf-8') as f:
                    lock_data = json.load(f)
                # Check if lock is stale
                if time.time() - lock_data.get('timestamp', 0) > 300:
                    self._release_lock()
                    return None
                return lock_data.get('user', 'Unknown')
            return None
        except Exception:
            return None
    
    def _get_skipped_files(self):
        """Get list of files that were skipped due to processing issues"""
        skipped_files = set()
        try:
            if self.log_file.exists():
                with open(self.log_file, 'r', encoding='utf-8') as f:
                    for line in f:
                        if "SKIPPED FILES:" in line:
                            # Read the next lines until we hit the separator
                            for next_line in f:
                                if next_line.startswith("  - "):
                                    file_path = next_line.strip().split(": ")[0][4:]  # Remove "  - " prefix
                                    skipped_files.add(file_path)
                                elif "=" in next_line:
                                    break
                            break
        except Exception:
            pass
        return skipped_files

def main():
    st.set_page_config(
        page_title="SustainaCube RAG",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    st.title("🌱 SustainaCube: Sustainability ExpertCenter")
    st.markdown("Ask questions about sustainability, recycling, and environmental research in the Polyurethane Industry")
    
    # Help / User Manual overlay
    with st.expander("📖 Open User Manual"):
        try:
            from pathlib import Path
            manual_path = Path(__file__).parent / "USER_MANUAL.md"
            if manual_path.exists():
                st.markdown(manual_path.read_text(encoding="utf-8"))
                if st.button("Open manual as HTML (for PDF export)"):
                    html_path = Path(__file__).parent / "USER_MANUAL.html"
                    if html_path.exists():
                        st.markdown(f"Manual HTML is available at: `{html_path}`. Open it in your browser and print to PDF.")
                    else:
                        st.info("Manual HTML not found. Run `export_user_manual.bat` to generate it.")
            else:
                st.info("User manual not found. See `SustainaCube_RAG/USER_MANUAL.md`.")
        except Exception as _e:
            st.info("User manual could not be displayed.")
    
    # Initialize RAG system
    if 'rag_system' not in st.session_state:
        st.session_state.rag_system = SustainaCubeMinimal()
    
    # User identification and lock management
    st.session_state.user_name = st.selectbox(
        "👤 Select your name:", 
        options=["Stefan Hermes", "Bart ten Brink"],
        index=0 if st.session_state.get('user_name') != "Bart ten Brink" else 1,
        key="user_name_select"
    )
    
    # Check who's processing documents (not just using the app)
    current_user = st.session_state.user_name
    processing_user = st.session_state.rag_system._check_who_online()
    
    # Show processing status only when someone is processing
    if processing_user and processing_user != current_user:
        st.warning(f"⚠️ {processing_user} is currently processing documents")
    elif processing_user == current_user:
        st.success(f"✅ You ({current_user}) are currently processing documents")
    
    # Sidebar for document processing and statistics
    with st.sidebar:
        st.header("📚 Document Management")
        
        # Check for new files first
        new_files_count = st.session_state.rag_system.check_for_new_files("../Document Database")
        
        # Calculate actual new files based on current stats
        if st.session_state.rag_system.processed:
            total_available = len(st.session_state.rag_system._get_document_files("../Document Database"))
            processed_files = len({d.get('file_path') for d in st.session_state.rag_system.documents if d.get('file_path')})
            # Get skipped files from log to exclude them from "new files" count
            skipped_files = st.session_state.rag_system._get_skipped_files()
            # Exclude both successfully processed files AND skipped files from "new files" count
            actual_new_files = max(0, total_available - processed_files - len(skipped_files))
        else:
            actual_new_files = new_files_count
        
        if actual_new_files > 0:
            # Show new documents available indicator
            st.info(f"🆕 {actual_new_files} new or updated files detected!")
            
            # Check if someone else is processing
            if processing_user and processing_user != current_user:
                st.error(f"❌ Cannot process documents - {processing_user} is currently processing documents")
            else:
                # Active load button with green color
                if st.button("🔄 Load New Documents", type="primary"):
                    # Acquire lock before processing
                    lock_acquired, lock_user = st.session_state.rag_system._acquire_lock(current_user)
                    if not lock_acquired:
                        st.error(f"❌ Cannot acquire lock - {lock_user} is processing documents")
                    else:
                        try:
                            with st.spinner(f"Loading new documents... ({current_user} is processing)"):
                                st.session_state.rag_system.process_documents("../Document Database")
                                
                                # Show completion alert
                                if st.session_state.rag_system.processed:
                                    st.success(f"✅ Successfully processed {len({d.get('file_path') for d in st.session_state.rag_system.documents})} files!")
                                    st.rerun()  # Refresh to update the UI
                        finally:
                            # Release lock after processing
                            st.session_state.rag_system._release_lock()
        else:
            # Disabled button when no new files
            st.button("🔄 Load Documents", disabled=True, help="No new documents to load")
            st.caption("All documents are up to date")
        
        st.markdown("---")
        if st.session_state.rag_system.processed:
            st.success("✅ Knowledge base ready!")
        else:
            st.warning("⚠️ Process documents first")
        
        # Statistics section
        st.markdown("### 📊 Quick Stats")
        
        # Get total files available in Document Database
        doc_files = st.session_state.rag_system._get_document_files("../Document Database")
        total_available = len(doc_files)
        
        # Debug: show file type breakdown
        if total_available > 0:
            from collections import Counter
            extensions = [Path(f).suffix.lower() for f in doc_files]
            ext_counts = Counter(extensions)
            st.caption(f"File types: {dict(ext_counts)}")
        
        # Show processing status in sidebar
        if processing_user and processing_user != current_user:
            st.warning(f"⚠️ {processing_user} is currently processing documents")
        elif processing_user == current_user:
            st.success(f"✅ You ({current_user}) are currently processing documents")
        else:
            st.info("ℹ️ No one is currently processing documents")
        
        # Show unprocessed files
        if st.session_state.rag_system.processed:
            processed_files = {d.get('file_path') for d in st.session_state.rag_system.documents if d.get('file_path')}
            unprocessed_files = [f for f in doc_files if f not in processed_files]
            if unprocessed_files:
                st.markdown("### 📋 Unprocessed Files")
                if st.button("📄 Show Unprocessed Files"):
                    st.text_area("Files not yet processed:", "\n".join(unprocessed_files), height=200)
        
        # Log file access button
        st.markdown("### 📋 Processing Log")
        col1, col2 = st.columns(2)
        with col1:
            if st.button("📄 View Processing Log"):
                if st.session_state.rag_system.log_file.exists():
                    with open(st.session_state.rag_system.log_file, "r", encoding="utf-8") as f:
                        log_content = f.read()
                    st.text_area("Processing Log", log_content, height=300)
                else:
                    st.info("No log file found yet. Process some documents first.")
        with col2:
            if st.button("🗑️ Clear Log Files", help="Clear only log files (keeps knowledge base intact)"):
                if st.session_state.rag_system.log_file.exists():
                    st.session_state.rag_system.log_file.unlink()
                if st.session_state.rag_system.progress_path.exists():
                    st.session_state.rag_system.progress_path.unlink()
                if st.session_state.rag_system.lock_file.exists():
                    st.session_state.rag_system.lock_file.unlink()
                st.success("✅ Log files cleared! (Knowledge base preserved)")
                st.rerun()
        
        if st.session_state.rag_system.processed:
            # Count unique files that have been processed
            processed_files = len({d.get('file_path') for d in st.session_state.rag_system.documents if d.get('file_path')})
            st.metric("Total Files Available", total_available)
            st.metric("Files Processed", processed_files)
            if total_available > processed_files:
                st.metric("New Files", total_available - processed_files)
            else:
                st.metric("New Files", 0)
        else:
            st.metric("Total Files Available", total_available)
            st.metric("Files Processed", 0)
            if total_available > 0:
                st.metric("New Files", total_available)
    
    # Main interface
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.header("💬 Ask a Question")
        
        # Manage question state
        if 'question_input' not in st.session_state:
            st.session_state.question_input = ""
        if 'auto_run' not in st.session_state:
            st.session_state.auto_run = False

        # Question input
        question = st.text_area(
            "Enter your sustainability question:",
            value=st.session_state.question_input,
            placeholder="e.g., What are the CO2 savings from PU foam recycling in Thailand?",
            height=100
        )
        st.session_state.question_input = question

        def run_query(q: str):
            with st.spinner("Searching knowledge base and generating answer..."):
                answer, sources = st.session_state.rag_system.answer_question(q)
            st.markdown("### 📋 Answer")
            st.markdown(answer)
            # Convert markdown to HTML and provide professional styling
            import re
            
            # Convert markdown headers to HTML
            html_answer = answer
            html_answer = re.sub(r'^### (.+)$', r'<h3>\1</h3>', html_answer, flags=re.MULTILINE)
            html_answer = re.sub(r'^#### (.+)$', r'<h4>\1</h4>', html_answer, flags=re.MULTILINE)
            html_answer = re.sub(r'^## (.+)$', r'<h2>\1</h2>', html_answer, flags=re.MULTILINE)
            html_answer = re.sub(r'^# (.+)$', r'<h1>\1</h1>', html_answer, flags=re.MULTILINE)
            
            # Remove source references from text (keep only the clean content)
            html_answer = re.sub(r'【[^】]+】', '', html_answer)
            
            # Remove "Source References" section from content since we have a separate sources section
            html_answer = re.sub(r'<h3>Source References</h3>.*?(?=<h3>|$)', '', html_answer, flags=re.DOTALL)
            
            # Convert bullet points
            html_answer = re.sub(r'^- (.+)$', r'<li>\1</li>', html_answer, flags=re.MULTILINE)
            # Wrap consecutive <li> in <ul>
            html_answer = re.sub(r'(<li>.*</li>)(?:\s*<li>.*</li>)*', lambda m: f'<ul>{m.group(0)}</ul>', html_answer, flags=re.DOTALL)
            
            # Convert bold text
            html_answer = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', html_answer)
            
            # Convert line breaks
            html_answer = html_answer.replace('\n', '<br>')
            
            # Extract and format source references for the bottom
            source_refs = []
            if sources:
                for source in sources:
                    src_name = source.get('filename') if isinstance(source, dict) else str(source)
                    score = source.get('similarity_score') if isinstance(source, dict) else None
                    if isinstance(score, (int, float)):
                        source_refs.append(f'<li>{src_name} (Relevance: {score:.3f})</li>')
                    else:
                        source_refs.append(f'<li>{src_name}</li>')
            source_list = f'<ul>{"".join(source_refs)}</ul>' if source_refs else '<p><em>No specific sources referenced.</em></p>'
            
            html_content = f"""
<!doctype html>
<html>
<head>
  <meta charset=\"utf-8\" />
  <meta name=\"viewport\" content=\"width=device-width, initial-scale=1\" />
  <style>
    body {{ 
      font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, 'Helvetica Neue', Arial, sans-serif; 
      line-height: 1.7; 
      color: #2c3e50; 
      margin: 0; 
      padding: 40px; 
      background: #f8f9fa;
    }}
    .container {{
      max-width: 800px; 
      margin: 0 auto; 
      background: white; 
      padding: 40px; 
      border-radius: 8px; 
      box-shadow: 0 2px 10px rgba(0,0,0,0.1);
    }}
    h1 {{ 
      color: #2c3e50; 
      font-size: 28px; 
      font-weight: 700; 
      margin: 0 0 20px 0; 
      border-bottom: 3px solid #3498db; 
      padding-bottom: 10px;
    }}
    h2 {{ 
      color: #34495e; 
      font-size: 22px; 
      font-weight: 600; 
      margin: 30px 0 15px 0; 
      border-left: 4px solid #3498db; 
      padding-left: 15px;
    }}
    h3 {{ 
      color: #34495e; 
      font-size: 18px; 
      font-weight: 600; 
      margin: 25px 0 12px 0;
    }}
    h4 {{ 
      color: #34495e; 
      font-size: 16px; 
      font-weight: 600; 
      margin: 20px 0 10px 0;
    }}
    ul {{ 
      margin: 15px 0; 
      padding-left: 20px;
    }}
    li {{ 
      margin: 8px 0; 
      line-height: 1.6;
    }}
    p {{ 
      margin: 15px 0; 
      line-height: 1.7;
    }}
    strong {{ 
      color: #2c3e50; 
      font-weight: 600;
    }}
    .header {{
      text-align: center; 
      margin-bottom: 30px; 
      padding-bottom: 20px; 
      border-bottom: 2px solid #ecf0f1;
    }}
    .header h1 {{ 
      border: none; 
      margin: 0; 
      color: #2c3e50;
    }}
    .timestamp {{
      color: #7f8c8d; 
      font-size: 14px; 
      margin-top: 10px;
    }}
    .question-section {{
      background: #f8f9fa; 
      padding: 20px; 
      border-radius: 6px; 
      margin: 20px 0; 
      border-left: 4px solid #3498db;
    }}
    .question-text {{
      font-size: 16px; 
      color: #2c3e50; 
      font-weight: 500; 
      margin: 10px 0 0 0; 
      line-height: 1.6;
    }}
    .sources-section {{
      margin-top: 30px; 
      padding-top: 20px; 
      border-top: 2px solid #ecf0f1;
    }}
    .sources-section h2 {{
      color: #7f8c8d; 
      font-size: 16px; 
      margin-bottom: 15px;
    }}
  </style>
  <title>SustainaCube Expert Response</title>
</head>
<body>
  <div class=\"container\">
    <div class=\"header\">
      <h1>🌱 SustainaCube Expert Response</h1>
      <div class=\"timestamp\">Generated on {time.strftime('%B %d, %Y at %I:%M %p')}</div>
    </div>
    <div class="question-section">
      <h2>❓ Question</h2>
      <p class="question-text">{q}</p>
    </div>
    <div class=\"content\">{html_answer}</div>
    <div class="sources-section">
      <h2>📚 Sources</h2>
      {source_list}
    </div>
  </div>
</body>
</html>
"""
            st.download_button(
                label="Copy Answer (HTML)",
                data=html_content,
                file_name="sustainacube_answer.html",
                mime="text/html",
                key="download_html"
            )
            st.download_button(
                label="Copy Answer (Text)",
                data=answer,
                file_name="sustainacube_answer.txt",
                mime="text/plain",
                key="download_text"
            )
            if sources:
                st.markdown("### 📚 Sources")
                for source in sources:
                    name = source.get('filename') if isinstance(source, dict) else str(source)
                    score = source.get('similarity_score') if isinstance(source, dict) else None
                    line = f"- **{name}**"
                    if isinstance(score, (int, float)):
                        line += f" (Relevance: {score:.3f})"
                    st.markdown(line)
                    matched = source.get('matched_words') if isinstance(source, dict) else None
                    if matched:
                        st.markdown(f"  *Matched: {', '.join(matched[:5])}...*")

        if st.button("🔍 Get Answer", type="primary"):
            if question.strip():
                run_query(question)
            else:
                st.warning("Please enter a question.")

        # Auto-run if triggered by a sample click
        if st.session_state.auto_run and st.session_state.question_input.strip():
            run_query(st.session_state.question_input)
            st.session_state.auto_run = False
    
    with col2:
        # Right column content can be added here if needed
        pass
        
        st.markdown("---")
        st.markdown("### 🤖 Assistant")
        use_assistant = st.checkbox("Use OpenAI Assistant (Vector Store)", value=bool(st.session_state.rag_system.assistant_id))
        # Hide Assistant ID; use environment-configured value only
        if use_assistant:
            # keep current assistant_id (from .env) if present
            if not st.session_state.rag_system.assistant_id:
                st.info("Assistant ID not found in environment. Set OPENAI_ASSISTANT_ID in .env and restart.")
            else:
                st.caption("Assistant enabled. Answers will use your OpenAI Assistant with Retrieval/WebSearch if enabled.")
        else:
            st.session_state.rag_system.assistant_id = ""
        st.markdown("---")
        st.markdown("### 💡 Sample Questions")
        sample_questions = [
            "What are the environmental benefits of PU foam recycling?",
            "Compare EPR frameworks across different countries",
            "What are the latest chemical recycling methods?",
            "How much CO2 can be saved through mattress recycling?",
            "What are the economic benefits of circular economy?"
        ]
        # Clickable buttons that fill input and auto-run
        for q in sample_questions:
            if st.button(q, key=f"sample_btn_{q}"):
                st.session_state.question_input = q
                st.session_state.auto_run = True
                st.rerun()

        # Copyable text list
        st.markdown("### 📋 Copyable Samples")
        st.code("\n".join(f"- {q}" for q in sample_questions))

if __name__ == "__main__":
    main()
