import streamlit as st
# Build: simple-sheets integration active
import openai
import os
from pathlib import Path
import PyPDF2
from docx import Document
from dotenv import load_dotenv
import re
from collections import Counter
import json
import time
from datetime import datetime
from PIL import Image
from typing import Dict

# Load environment variables
load_dotenv()

class SustainaCubeMinimal:
    def __init__(self):
        # Try Streamlit secrets first, fallback to environment variables
        try:
            api_key = st.secrets["OPENAI_API_KEY"]
            assistant_id = st.secrets.get("OPENAI_ASSISTANT_ID", "")
        except:
            # Fallback to environment variables for local development
            api_key = os.getenv("OPENAI_API_KEY")
            assistant_id = os.getenv("OPENAI_ASSISTANT_ID", "")
        
        self.openai_client = openai.OpenAI(api_key=api_key)
        self.documents = []
        self.processed = False
        # Assistant (OpenAI Assistants API)
        self.assistant_id = assistant_id
        # Resume/cache/tracking
        self.cache_path = Path(__file__).parent / "processed_cache.json"
        self.chunks_store_path = Path(__file__).parent / "chunks_store.jsonl"
        self.progress_path = Path(__file__).parent / ".progress.json"
        self.log_file = Path(__file__).parent / "processing_log.txt"
        self.log_entries = []
        self.lock_file = Path(__file__).parent / "app_lock.json"
        self._processed_index = {}  # file_path -> {mtime, size}
        self._load_cache_and_chunks()
        
    def extract_text_from_file(self, file_path):
        """Extract text from various file formats"""
        try:
            file_path = Path(file_path)
            if file_path.suffix.lower() == '.pdf':
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                    
                    # If no text extracted, try OCR
                    if not text.strip():
                        try:
                            import pytesseract
                            from PIL import Image
                            import fitz  # PyMuPDF for better PDF handling
                            
                            doc = fitz.open(file_path)
                            text = ""
                            for page_num in range(len(doc)):
                                page = doc.load_page(page_num)
                                pix = page.get_pixmap()
                                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                                text += pytesseract.image_to_string(img) + "\n"
                            doc.close()
                        except ImportError:
                            pass
                    
                    return text
            elif file_path.suffix.lower() in ['.docx', '.doc']:
                doc = Document(file_path)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                return text
            elif file_path.suffix.lower() == '.txt':
                with open(file_path, 'r', encoding='utf-8') as file:
                    return file.read()
            elif file_path.suffix.lower() == '.pptx':
                from pptx import Presentation
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            elif file_path.suffix.lower() in ['.xlsx', '.xls']:
                import openpyxl
                wb = openpyxl.load_workbook(file_path)
                text = ""
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    for row in sheet.iter_rows():
                        for cell in row:
                            if cell.value:
                                text += str(cell.value) + " "
                        text += "\n"
                return text
            elif file_path.suffix.lower() == '.html':
                with open(file_path, 'r', encoding='utf-8') as file:
                    content = file.read()
                    # Simple HTML tag removal
                    import re
                    text = re.sub(r'<[^>]+>', '', content)
                    return text
            else:
                return ""
        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
            return ""

    def chunk_text(self, text, chunk_size=1000, overlap=200):
        """Split text into overlapping chunks"""
        if not text.strip():
            return []
        
        words = text.split()
        chunks = []
        
        for i in range(0, len(words), chunk_size - overlap):
            chunk = " ".join(words[i:i + chunk_size])
            if chunk.strip():
                chunks.append(chunk)
        
        return chunks

    def process_documents(self, documents_folder):
        """Process all documents in the folder"""
        documents_folder = Path(documents_folder)
        if not documents_folder.exists():
            print(f"Documents folder not found: {documents_folder}")
            return
        
        self.log_entries = []
        self._log("Starting document processing...")
        
        # Get all document files
        doc_files = self._get_document_files(documents_folder)
        self._log(f"Found {len(doc_files)} files to process")
        
        # Process each file
        for i, file_path in enumerate(doc_files):
            try:
                self._log(f"Processing {i+1}/{len(doc_files)}: {file_path.name}")
                
                # Check if file was already processed and hasn't changed
                file_stat = file_path.stat()
                file_key = str(file_path)
                
                if file_key in self._processed_index:
                    cached_stat = self._processed_index[file_key]
                    if (cached_stat['mtime'] == file_stat.st_mtime and 
                        cached_stat['size'] == file_stat.st_size):
                        self._log(f"Skipping unchanged file: {file_path.name}")
                        continue
                
                # Extract text
                text = self.extract_text_from_file(file_path)
                if not text.strip():
                    self._log(f"No text extracted from: {file_path.name}")
                    continue
                
                # Chunk the text
                chunks = self.chunk_text(text)
                self._log(f"Created {len(chunks)} chunks from {file_path.name}")
                
                # Store chunks
                for j, chunk in enumerate(chunks):
                    doc_entry = {
                        'file_path': str(file_path),
                        'filename': file_path.name,
                        'chunk_id': f"{file_path.stem}_{j}",
                        'text': chunk,
                        'processed_at': datetime.now().isoformat()
                    }
                    self.documents.append(doc_entry)
                
                # Update processed index
                self._processed_index[file_key] = {
                    'mtime': file_stat.st_mtime,
                    'size': file_stat.st_size
                }
                
                # Save progress
                self._save_progress(i + 1, len(doc_files))
                
            except Exception as e:
                self._log(f"Error processing {file_path.name}: {e}")
                continue
        
        # Save all data
        self._save_documents()
        self._save_chunks()
        self._save_cache()
        self._save_log()
        
        self.processed = True
        self._log(f"Document processing completed. Processed {len(self.documents)} chunks from {len(doc_files)} files.")

    def _get_document_files(self, documents_folder):
        """Get all document files from the folder"""
        documents_folder = Path(documents_folder)
        extensions = ['.pdf', '.docx', '.doc', '.txt', '.pptx', '.xlsx', '.xls', '.html']
        files = []
        
        for ext in extensions:
            files.extend(documents_folder.glob(f"**/*{ext}"))
        
        return sorted(files)

    def _log(self, message):
        """Log a message with timestamp"""
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        log_entry = f"[{timestamp}] {message}"
        self.log_entries.append(log_entry)
        print(log_entry)

    def _save_progress(self, current, total):
        """Save progress to file"""
        progress = {
            'current': current,
            'total': total,
            'timestamp': datetime.now().isoformat()
        }
        with open(self.progress_path, 'w') as f:
            json.dump(progress, f)

    def _save_documents(self):
        """Save processed documents to cache"""
        with open(self.cache_path, 'w') as f:
            json.dump(self.documents, f, indent=2)

    def _save_chunks(self):
        """Save chunks to JSONL file"""
        with open(self.chunks_store_path, 'w') as f:
            for doc in self.documents:
                f.write(json.dumps(doc) + '\n')

    def _save_cache(self):
        """Save processed index cache"""
        cache_data = {
            'processed_index': self._processed_index,
            'last_updated': datetime.now().isoformat()
        }
        with open(self.cache_path.parent / "processed_index.json", 'w') as f:
            json.dump(cache_data, f, indent=2)

    def _save_log(self):
        """Save log entries to file"""
        with open(self.log_file, 'w', encoding='utf-8') as f:
            f.write('\n'.join(self.log_entries))

    def _load_cache_and_chunks(self):
        """Load cached documents and chunks"""
        try:
            if self.cache_path.exists():
                with open(self.cache_path, 'r') as f:
                    self.documents = json.load(f)
                self.processed = True
                self._log(f"Loaded {len(self.documents)} cached documents")
        except Exception as e:
            self._log(f"Error loading cache: {e}")

    def search_documents(self, query, top_k=5):
        """Search through processed documents"""
        if not self.documents:
            return []
        
        query_lower = query.lower()
        results = []
        
        for doc in self.documents:
            text_lower = doc['text'].lower()
            
            # Simple keyword matching
            matches = 0
            matched_words = []
            for word in query_lower.split():
                if word in text_lower:
                    matches += 1
                    matched_words.append(word)
            
            if matches > 0:
                # Calculate similarity score
                similarity = matches / len(query.split())
                results.append({
                    'filename': doc['filename'],
                    'text': doc['text'],
                    'similarity': similarity,
                    'matched_words': matched_words
                })
        
        # Sort by similarity and return top results
        results.sort(key=lambda x: x['similarity'], reverse=True)
        return results[:top_k]

    def answer_question(self, question):
        """Main method to answer a question using Assistant if configured, else local retrieval"""
        # Prefer Assistant (Vector Store + WebSearch) when configured
        if self.assistant_id:
            try:
                thread = self.openai_client.beta.threads.create()
                self.openai_client.beta.threads.messages.create(
                    thread_id=thread.id,
                    role="user",
                    content=question
                )
                run = self.openai_client.beta.threads.runs.create(
                    thread_id=thread.id,
                    assistant_id=self.assistant_id
                )
                # Poll until completion
                start = time.time()
                while True:
                    r = self.openai_client.beta.threads.runs.retrieve(thread_id=thread.id, run_id=run.id)
                    if r.status in ["completed", "failed", "cancelled", "expired"]:
                        break
                    time.sleep(0.4)
                    # optional safety timeout
                    if time.time() - start > 120:
                        return "Assistant timeout. Please try again.", []
                if r.status != "completed":
                    return f"Assistant run status: {r.status}", []
                msgs = self.openai_client.beta.threads.messages.list(thread_id=thread.id)
                answer_text = ""
                sources = []
                for m in msgs.data:
                    if m.role == 'assistant':
                        parts = []
                        for c in m.content:
                            if getattr(c, 'type', '') == 'text':
                                parts.append(c.text.value)
                        if parts:
                            answer_text = "\n\n".join(parts)
                        # If the assistant returns file references as attachments
                        if hasattr(m, 'attachments') and m.attachments:
                            for att in m.attachments:
                                fname = getattr(att, 'filename', None)
                                if fname:
                                    sources.append({'filename': fname, 'similarity_score': 1.0})
                        break
                
                # Extract source references from the answer text if no attachments found
                if not sources and answer_text:
                    import re
                    # Look for file references in the text (common patterns)
                    file_refs = re.findall(r'\[([^\]]+\.(?:pdf|docx?|txt|pptx?|xlsx?|html?|md|csv))\]', answer_text, re.IGNORECASE)
                    for ref in file_refs:
                        sources.append({'filename': ref, 'similarity_score': 1.0})
                    
                    # Also look for quoted filenames
                    quoted_files = re.findall(r'"([^"]+\.(?:pdf|docx?|txt|pptx?|xlsx?|html?|md|csv))"', answer_text, re.IGNORECASE)
                    for ref in quoted_files:
                        sources.append({'filename': ref, 'similarity_score': 1.0})
                
                return (answer_text or ""), sources
            except Exception as e:
                return f"Assistant error: {e}", []
        
        # Fallback to local retrieval
        search_results = self.search_documents(question)
        if not search_results:
            return "No relevant information found in the knowledge base.", []
        answer = self.generate_answer(question, search_results)
        sources = []
        for result in search_results:
            sources.append({'filename': result['filename'], 'similarity_score': result['similarity'], 'matched_words': result['matched_words']})
        return answer, sources

    def generate_answer(self, question, search_results):
        """Generate an answer based on search results"""
        # Limit context to avoid token limit (roughly 3000 characters max)
        context_parts = []
        total_length = 0
        max_length = 3000
        
        for result in search_results:
            text = result['text']
            if total_length + len(text) < max_length:
                context_parts.append(text)
                total_length += len(text)
            else:
                # Add partial text if there's room
                remaining = max_length - total_length
                if remaining > 100:  # Only add if there's meaningful space
                    context_parts.append(text[:remaining] + "...")
                break
        
        context = "\n\n".join(context_parts)
        
        prompt = f"""Based on the following context about sustainability in the polyurethane industry, answer the question: {question}

Context:
{context}

Answer:"""
        
        try:
            response = self.openai_client.chat.completions.create(
                model="gpt-4",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=1000,
                temperature=0.7
            )
            return response.choices[0].message.content
        except Exception as e:
            return f"Error generating answer: {e}"

    def check_for_new_files(self, documents_folder):
        """Check if there are new files to process"""
        doc_files = self._get_document_files(documents_folder)
        return len(doc_files)

    def _get_skipped_files(self):
        """Get list of skipped files from log"""
        skipped = []
        if self.log_file.exists():
            with open(self.log_file, 'r', encoding='utf-8') as f:
                for line in f:
                    if "No text extracted from:" in line:
                        filename = line.split("No text extracted from: ")[-1].strip()
                        skipped.append(filename)
        return skipped

    def _acquire_lock(self, user):
        """Acquire processing lock"""
        if self.lock_file.exists():
            with open(self.lock_file, 'r') as f:
                lock_data = json.load(f)
            return False, lock_data.get('user', 'unknown')
        
        lock_data = {'user': user, 'timestamp': datetime.now().isoformat()}
        with open(self.lock_file, 'w') as f:
            json.dump(lock_data, f)
        return True, user

    def _release_lock(self):
        """Release processing lock"""
        if self.lock_file.exists():
            self.lock_file.unlink()

def check_password():
    """User authentication using local CSV (offline-managed)"""
    if 'authenticated' not in st.session_state:
        st.session_state.authenticated = False

    # Load users from local CSV (prefer 'Corporate Users.csv', fallback to 'corporate_users.csv')
    import csv
    from pathlib import Path
    users: Dict[str, Dict] = {}
    base_dir = Path(__file__).parent
    csv_paths = [
        base_dir / "Corporate Users.csv",
        base_dir / "corporate_users.csv",
        Path("Corporate Users.csv"),
        Path("corporate_users.csv"),
    ]
    csv_loaded = False
    for csv_path in csv_paths:
        try:
            with open(csv_path, newline="", encoding="utf-8") as f:
                reader = csv.DictReader(f)
                for row in reader:
                    email_key = (row.get("Email", "") or "").strip().lower()
                    if email_key:
                        users[email_key] = {
                            'email': (row.get('Email', '') or '').strip(),
                            'password': (row.get('Password', '') or '').strip(),
                            'valid_until': (row.get('Valid_Until', '') or '').strip(),
                            'status': (row.get('Status', 'Active') or 'Active').strip(),
                        }
                csv_loaded = True
                break
        except FileNotFoundError:
            continue
    if not csv_loaded:
        st.error("Corporate Users.csv not found. Place it next to app_corporate.py or at repo root and redeploy.")
        return False

    if not st.session_state.authenticated:
        col1, col2 = st.columns([1, 4])
        with col1:
            try:
                st.image("Logo Carpe Diem 5.png", width=120)
            except:
                st.markdown("🌱")
        with col2:
            st.title("🔐 Corporate Access")

        st.markdown("Please enter your corporate credentials to access SustainaCube.")

        email = st.text_input("Email Address", placeholder="user@company.com")
        password = st.text_input("Password", type="password")

        if st.button("Login"):
            if email and password:
                user_id = email.lower().strip()
                if user_id in users:
                    user = users[user_id]
                    if user.get('status', 'Active').lower() != 'active':
                        st.error("Your account is not active. Contact administrator.")
                    elif user.get('password') == password:
                        try:
                            valid_until_str = user.get('valid_until', '')
                            if valid_until_str:
                                valid_until = datetime.strptime(valid_until_str, '%d/%m/%Y').date()
                                if valid_until < datetime.now().date():
                                    st.error("Your account has expired. Contact administrator.")
                                    return False
                            st.session_state.authenticated = True
                            st.session_state.current_user = user['email']
                            st.rerun()
                        except Exception:
                            st.session_state.authenticated = True
                            st.session_state.current_user = user['email']
                            st.rerun()
                    else:
                        st.error("Incorrect password. Please try again.")
                else:
                    st.error("User not found. Please contact your administrator.")
            else:
                st.error("Please enter both email and password.")
        return False
    return True

def main():
    # Check authentication first
    if not check_password():
        return
    
    st.set_page_config(
        page_title="SustainaCube - Corporate Version",
        layout="wide",
        initial_sidebar_state="collapsed"  # Hide sidebar for corporate version
    )
    
    # Header with logo
    col1, col2 = st.columns([1, 4])
    with col1:
        try:
            st.image("Logo Carpe Diem 5.png", width=120)
        except:
            st.markdown("🌱")  # Fallback if logo not found
    with col2:
        st.title("SustainaCube - Corporate Version")
    
    # Top bar: right-aligned logout button (no "logged in as" text)
    top_left, top_right = st.columns([6, 1])
    with top_right:
        if st.button("🚪 Logout", use_container_width=True):
            st.session_state.authenticated = False
            if 'current_user' in st.session_state:
                del st.session_state.current_user
            st.rerun()

    st.markdown("Ask questions about sustainability, recycling, and environmental research in the Polyurethane Industry.")
    
    # User Manual section
    with st.expander("📖 Open User Manual"):
        try:
            from pathlib import Path
            manual_path = Path(__file__).parent / "USER_MANUAL_CORPORATE.md"
            if manual_path.exists():
                st.markdown(manual_path.read_text(encoding="utf-8"))
            else:
                st.info("User manual not found. See `USER_MANUAL_CORPORATE.md` in the repository.")
        except Exception as _e:
            st.info("User manual could not be displayed.")
    
    # Initialize RAG system
    if 'rag_system' not in st.session_state:
        st.session_state.rag_system = SustainaCubeMinimal()
    
    # Process documents automatically on first load
    if not st.session_state.rag_system.processed:
        with st.spinner("Loading knowledge base..."):
            st.session_state.rag_system.process_documents("./Document Database")
    
    # Main interface
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.header("💬 Ask a Question")
        
        # Manage question state
        if 'question_input' not in st.session_state:
            st.session_state.question_input = ""
        if 'auto_run' not in st.session_state:
            st.session_state.auto_run = False

        # Question input
        question = st.text_area(
            "Enter your sustainability question:",
            value=st.session_state.question_input,
            placeholder="e.g., What are the CO2 savings from PU foam recycling in Thailand?",
            height=100
        )
        st.session_state.question_input = question

        def run_query(q: str):
            with st.spinner("Searching knowledge base and generating answer..."):
                answer, sources = st.session_state.rag_system.answer_question(q)
            
            # Usage tracking removed - using CSV authentication only
            
            st.markdown("### 📋 Answer")
            st.markdown(answer)
            
            # Convert markdown to HTML and provide professional styling
            import re
            import time
            
            # Convert markdown headers to HTML
            html_answer = answer
            html_answer = re.sub(r'^### (.+)$', r'<h3>\1</h3>', html_answer, flags=re.MULTILINE)
            html_answer = re.sub(r'^#### (.+)$', r'<h4>\1</h4>', html_answer, flags=re.MULTILINE)
            html_answer = re.sub(r'^## (.+)$', r'<h2>\1</h2>', html_answer, flags=re.MULTILINE)
            html_answer = re.sub(r'^# (.+)$', r'<h1>\1</h1>', html_answer, flags=re.MULTILINE)
            
            # Remove all source references from text for corporate version
            html_answer = re.sub(r'【[^】]+】', '', html_answer)
            html_answer = re.sub(r'<h3>Source References</h3>.*?(?=<h3>|$)', '', html_answer, flags=re.DOTALL)
            html_answer = re.sub(r'Source References.*?(?=<h[1-6]|$)', '', html_answer, flags=re.DOTALL | re.IGNORECASE)
            html_answer = re.sub(r'References:.*?(?=<h[1-6]|$)', '', html_answer, flags=re.DOTALL | re.IGNORECASE)
            html_answer = re.sub(r'\[.*?\.(?:pdf|docx?|txt|pptx?|xlsx?|html?|md|csv)\]', '', html_answer, flags=re.IGNORECASE)
            
            # Convert bullet points
            html_answer = re.sub(r'^- (.+)$', r'<li>\1</li>', html_answer, flags=re.MULTILINE)
            # Wrap consecutive <li> in <ul>
            html_answer = re.sub(r'(<li>.*</li>)(?:\s*<li>.*</li>)*', lambda m: f'<ul>{m.group(0)}</ul>', html_answer, flags=re.DOTALL)
            
            # Convert bold text
            html_answer = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', html_answer)
            
            # Convert line breaks
            html_answer = html_answer.replace('\n', '<br>')
            
            # No source references for corporate version
            
            html_content = f"""
<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>SustainaCube - Corporate Response</title>
  <style>
    @page {{
      margin: 2cm;
      size: A4;
    }}
    body {{
      font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
      line-height: 1.5;
      color: #333;
      margin: 0;
      padding: 0;
      background: white;
    }}
    .container {{
      max-width: 21cm;
      margin: 0 auto;
      padding: 2cm;
      background: white;
    }}
    .header {{
      border-bottom: 3px solid #2c5aa0;
      padding-bottom: 20px;
      margin-bottom: 30px;
      display: flex;
      align-items: center;
      justify-content: space-between;
    }}
    .logo-section {{
      display: flex;
      align-items: center;
    }}
    .logo-text {{
      font-size: 28px;
      font-weight: bold;
      color: #2c5aa0;
      margin-left: 15px;
    }}
    .question-section {{
      background: #f8f9fa;
      padding: 15px;
      border-left: 4px solid #2c5aa0;
      margin-bottom: 25px;
      border-radius: 0 4px 4px 0;
    }}
    .question-label {{
      font-weight: bold;
      color: #2c5aa0;
      margin-bottom: 8px;
      font-size: 14px;
      text-transform: uppercase;
      letter-spacing: 0.5px;
    }}
    .question-text {{
      font-style: italic;
      color: #555;
      margin: 0;
      line-height: 1.4;
    }}
    .content {{
      font-size: 14px;
      line-height: 1.6;
    }}
    .content h1 {{
      color: #2c5aa0;
      font-size: 24px;
      margin: 25px 0 15px 0;
      border-bottom: 2px solid #e9ecef;
      padding-bottom: 8px;
    }}
    .content h2 {{
      color: #2c5aa0;
      font-size: 20px;
      margin: 20px 0 12px 0;
    }}
    .content h3 {{
      color: #2c5aa0;
      font-size: 16px;
      margin: 15px 0 10px 0;
    }}
    .content h4 {{
      color: #2c5aa0;
      font-size: 14px;
      margin: 12px 0 8px 0;
    }}
    .content ul {{
      margin: 10px 0;
      padding-left: 20px;
    }}
    .content li {{
      margin: 5px 0;
    }}
    .content strong {{
      color: #2c5aa0;
      font-weight: 600;
    }}
    .content p {{
      margin: 10px 0;
    }}
    .footer {{
      margin-top: 40px;
      padding-top: 20px;
      border-top: 1px solid #e9ecef;
      text-align: center;
      font-size: 12px;
      color: #666;
    }}
    .footer .branding {{
      font-weight: bold;
      color: #2c5aa0;
      margin-bottom: 5px;
    }}
    .footer .disclaimer {{
      font-style: italic;
      margin-top: 10px;
    }}
  </style>
</head>
<body>
  <div class="container">
    <div class="header">
      <div class="logo-section">
        <div class="logo-text">SustainaCube</div>
      </div>
      <div style="font-size: 12px; color: #666; text-align: right;">
        Corporate Version
      </div>
    </div>
    
    <div class="question-section">
      <div class="question-label">Question</div>
      <div class="question-text">{q}</div>
    </div>
    
    <div class="content">{html_answer}</div>
    
    <div class="footer">
      <div class="branding">Carpe Diem / HTC Global</div>
      <div>Generated on {datetime.now().strftime('%d %B %Y at %H:%M')}</div>
      <div class="disclaimer">Answers provided by the PU ExpertCenter</div>
      <div style="margin-top: 10px;">Copyright Carpe Diem/HTC Global, all rights reserved</div>
    </div>
  </div>
</body>
</html>
"""
            st.download_button(
                label="Copy Answer (HTML)",
                data=html_content,
                file_name="sustainacube_answer.html",
                mime="text/html",
                key="download_html"
            )
            st.download_button(
                label="Copy Answer (Text)",
                data=answer,
                file_name="sustainacube_answer.txt",
                mime="text/plain",
                key="download_text"
            )
            # Sources section removed for corporate version

        if st.button("🔍 Get Answer", type="primary"):
            if question.strip():
                run_query(question)
            else:
                st.warning("Please enter a question.")

        # Auto-run if triggered by a sample click
        if st.session_state.auto_run and st.session_state.question_input.strip():
            run_query(st.session_state.question_input)
            st.session_state.auto_run = False
    
    with col2:
        # Right column content
        st.markdown("### 💡 Sample Questions")
        sample_questions = [
            "What are the environmental benefits of PU foam recycling?",
            "Compare EPR frameworks across different countries",
            "What are the latest chemical recycling methods?",
            "How much CO2 can be saved through mattress recycling?",
            "What are the economic benefits of circular economy?"
        ]
        # Clickable buttons that fill input and auto-run
        for q in sample_questions:
            if st.button(q, key=f"sample_btn_{q}"):
                st.session_state.question_input = q
                st.session_state.auto_run = True
                st.rerun()

if __name__ == "__main__":
    main()
